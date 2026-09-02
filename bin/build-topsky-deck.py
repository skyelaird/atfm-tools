#!/usr/bin/env python
"""Assemble the TopSky induction PowerPoint deck.

Reads:
  - Slide content from this script's SLIDES list (Python-defined).
  - Screenshots from C:\\Users\\JoelMorin\\OneDrive\\Games\\VATSIM\\CTP\\topsky-induction\\screenshots\\
    referenced by filename (no path) per slide.

Writes:
  - C:\\Users\\JoelMorin\\OneDrive\\Games\\VATSIM\\CTP\\topsky-induction\\TopSky-Induction.pptx

Re-run any time content evolves or screenshots are updated.
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

sys.stdout.reconfigure(encoding='utf-8')

ROOT = r'C:\Users\JoelMorin\OneDrive\Games\VATSIM\CTP\topsky-induction'
SCREENSHOT_DIR = os.path.join(ROOT, 'screenshots')
OUT_PATH = os.path.join(ROOT, 'TopSky-Induction.pptx')

# Colours
DARK_BG = RGBColor(0x10, 0x14, 0x20)
ACCENT = RGBColor(0x5e, 0xe0, 0xba)
WARN = RGBColor(0xf5, 0xa0, 0x5b)
TEXT = RGBColor(0xe8, 0xec, 0xf4)
DIM = RGBColor(0x9a, 0xa4, 0xb9)

# 16:9 deck — 13.333 in × 7.5 in
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# --------------- helpers ---------------

def add_blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(blank_layout)
    # Dark background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=16, color=TEXT, indent_px=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        run = p.add_run()
        run.text = '•  ' + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = 'Calibri'
        p.space_after = Pt(6)
    return tb


def add_image(slide, x, y, w, h, filename):
    path = os.path.join(SCREENSHOT_DIR, filename)
    if not os.path.exists(path):
        # Placeholder rectangle if screenshot missing
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor(0x2a, 0x31, 0x42)
        ph.line.color.rgb = WARN
        tf = ph.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f'[ screenshot missing: {filename} ]'
        run.font.size = Pt(14)
        run.font.italic = True
        run.font.color.rgb = WARN
        return ph
    return slide.shapes.add_picture(path, x, y, w, h)


def add_section_divider(prs, section_num, section_title):
    slide = add_blank_slide(prs)
    add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
             f'Section {section_num}', size=28, bold=True, color=ACCENT)
    add_text(slide, Inches(1), Inches(3.4), Inches(11), Inches(1.5),
             section_title, size=44, bold=True, color=TEXT)
    return slide


def add_title_slide(prs, title, subtitle, footer):
    slide = add_blank_slide(prs)
    add_text(slide, Inches(1), Inches(2.4), Inches(11), Inches(1.2),
             title, size=44, bold=True, color=ACCENT)
    add_text(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.7),
             subtitle, size=24, color=TEXT)
    add_text(slide, Inches(1), Inches(6.6), Inches(11), Inches(0.5),
             footer, size=14, color=DIM)
    return slide


def add_content_slide(prs, title, body_lines=None, image=None, image_caption=None,
                      image_pos='right'):
    """Standard content slide with title at top, body bullets and optional image.

    image_pos: 'right' | 'full' | 'below' | 'none'
    """
    slide = add_blank_slide(prs)
    # Title
    add_text(slide, Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.7),
             title, size=28, bold=True, color=ACCENT)
    # Accent underline
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0.5), Inches(1.05),
                                       Inches(0.7), Inches(0.04))
    underline.fill.solid()
    underline.fill.fore_color.rgb = ACCENT
    underline.line.fill.background()

    if image_pos == 'full' and image:
        # Full-width image below title
        add_image(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5), image)
        if image_caption:
            add_text(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
                     image_caption, size=12, color=DIM, align=PP_ALIGN.CENTER)
    elif image_pos == 'right' and image:
        # Body left, image right
        if body_lines:
            add_bullets(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.8),
                        body_lines, size=16)
        add_image(slide, Inches(6.3), Inches(1.3), Inches(6.5), Inches(5.5), image)
        if image_caption:
            add_text(slide, Inches(6.3), Inches(6.85), Inches(6.5), Inches(0.4),
                     image_caption, size=11, color=DIM, align=PP_ALIGN.CENTER)
    elif image_pos == 'below' and image:
        if body_lines:
            add_bullets(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2),
                        body_lines, size=16)
        add_image(slide, Inches(2), Inches(3.5), Inches(9.3), Inches(3.5), image)
        if image_caption:
            add_text(slide, Inches(2), Inches(7.05), Inches(9.3), Inches(0.4),
                     image_caption, size=11, color=DIM, align=PP_ALIGN.CENTER)
    else:
        # No image, full body
        if body_lines:
            add_bullets(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8),
                        body_lines, size=18)
    return slide


# --------------- deck definition ---------------

SLIDES = [
    # Title slide
    {'kind': 'title',
     'title': 'TopSky for EuroScope',
     'subtitle': 'A controller induction',
     'footer': 'For new-to-TopSky controllers and students · Setup B (COOPANS)'},

    # ---------- Section 1: Introduction ----------
    {'kind': 'divider', 'num': 1, 'name': 'Introduction'},

    {'kind': 'content',
     'title': 'What is TopSky?',
     'body': [
         'A plug-in for EuroScope that adds modern ATC functionality.',
         'Modelled on COOPANS — the operational system used by NAV CANADA, NATS, IAA, Naviair, Croatia Control, and Austro Control.',
         'Adds: track labels with full data fields, conflict detection (MTCD), safety nets (STCA, MSAW, AIW, APW), CPDLC, coordination tools.',
         'Replaces the basic EuroScope display with a much richer operational picture.',
     ]},

    {'kind': 'content',
     'title': 'TopSky vs vanilla EuroScope',
     'body': [
         'Vanilla EuroScope: basic radar display, simple labels, no advanced safety nets.',
         'TopSky: rich labels (CFL/RFL/AHDG/SSR/etc. as menu-driven fields), MTCD conflict probe, real-world ATC ergonomics.',
         'EuroScope is still the host; TopSky lives as a plug-in inside it.',
         'You still load EuroScope, connect to VATSIM, and load your sector — TopSky enhances the display once loaded.',
     ]},

    {'kind': 'content',
     'title': 'Scope of this induction',
     'body': [
         'The TopSky workspace and air situation display.',
         'How to read and manipulate track labels — the core controller interaction.',
         'The Global Menu and what each top-level item does.',
         'Key operational windows.',
         'Safety nets, monitoring aids, and the MTCD conflict probe.',
         'Practical workflows you will use every session.',
         'Reference: TopSky for EuroScope — General (COOPANS), version 2.6.',
     ]},

    # ---------- Section 2: Air Situation Display ----------
    {'kind': 'divider', 'num': 2, 'name': 'Air Situation Display'},

    {'kind': 'content',
     'title': 'The default workspace',
     'body': [
         'Global Menu bar (top): your main entry to all TopSky settings and tools.',
         'Air Situation Display (centre): radar picture with sector boundaries, airports, and aircraft tracks.',
         'AMAN button (top-left): the Arrival Manager.',
         'CPDLC Current Message Window (top-right): incoming/outgoing data link messages.',
         'Lists (bottom): Startup, Sector Inbound, Sector Exit, Departure.',
     ],
     'image': '01-default-workspace.png',
     'image_pos': 'below',
     'image_caption': 'Default TopSky workspace, ZQM_CTR_TS sector.'},

    {'kind': 'content',
     'title': 'Track presentation — what you see for each aircraft',
     'body': [
         'Position symbol: aircraft current location (varies by track type — controlled, assumed, etc.).',
         'History dots: trail of recent positions.',
         'Prediction line: a short forward vector showing where the aircraft will be in N seconds.',
         'Track label: the data block beside the position symbol — callsign, altitude, ground speed, route info.',
         'Colour coding: tells you the track ownership and state at a glance.',
     ]},

    {'kind': 'content',
     'title': 'Track label anatomy',
     'body': [
         'A multi-line, multi-field block showing the aircraft\'s controlled state.',
         'Each field is clickable — clicking opens a context menu specific to that field.',
         'Common fields: callsign, AFL (actual flight level), CFL (cleared), RFL (requested), AHDG (heading), SSR code, ARC (arrival), ASP (assigned speed).',
         'Section 3 covers each field menu in detail.',
     ]},

    {'kind': 'content',
     'title': 'Track filtering',
     'body': [
         'Altitude filter: show only tracks within a band.',
         'SSR code filter: show specific transponder codes.',
         'Quick Look: temporary display of tracks normally outside your area.',
         'Level Band Highlight: visually distinguish a chosen FL band.',
         'Oceanic Level highlight: special display for oceanic levels (relevant for CZQM/CZQX).',
         'All accessible via the Radar Menu (right-click on the display background).',
     ]},

    # ---------- Section 4: Global Menu (out of doc order; we put before Track Labels for tour flow) ----------
    {'kind': 'divider', 'num': 3, 'name': 'The Global Menu'},

    {'kind': 'content',
     'title': 'Setup menu',
     'body': [
         'Direction Finder, Weather Map, Unit Settings.',
         'Local Settings — tweak per-position preferences.',
         'CFL submenu, Flight Leg, RR Main.',
         'Brightness Control — adjust display brightness.',
         'CPDLC Setting — data link parameters.',
         'FAST and Blind Spot toggles.',
     ],
     'image': '10-global-setup-menu.png',
     'image_pos': 'right',
     'image_caption': 'Setup menu opened.'},

    {'kind': 'content',
     'title': 'AMS — Airspace Management',
     'body': [
         'FSA: Free Sectorisation Application — re-sectorise on the fly.',
         'NAT: North Atlantic Track management — particularly relevant for CZQM/CZQX oceanic operations.',
         'A small menu but important for FIR-level coordination.',
     ],
     'image': '11-global-ams-menu.png',
     'image_pos': 'right',
     'image_caption': 'AMS menu — FSA and NAT submenus.'},

    {'kind': 'content',
     'title': 'FData — Flight Plan data',
     'body': [
         'Flight Plan Selection: pick a flight to inspect.',
         'Flight Plan Window: show the full FPL detail.',
         'Used routinely when checking routes, EOBT, RFL.',
     ],
     'image': '12-global-fdata-menu.png',
     'image_pos': 'right',
     'image_caption': 'FData menu — Flight Plan tools.'},

    {'kind': 'content',
     'title': 'Tools — operational utilities',
     'body': [
         'Flight Plan Lists: managed views of flights by phase.',
         'CARD: Conflict And Risk Display — visualises MTCD conflicts.',
         'SAP: Segregated Area Probe.',
         'Message In/Out: ATC text messages.',
         'Shortcut: define position-specific shortcut buttons.',
         'CPDLC: data link operations.',
     ],
     'image': '13-global-tools-menu.png',
     'image_pos': 'right',
     'image_caption': 'Tools menu — operational utilities.'},

    {'kind': 'content',
     'title': 'MET — meteorology',
     'body': [
         'Messages: METAR/TAF/SIGMET/AIRMET retrieval.',
         'Upper Winds: winds aloft data for trajectory awareness.',
         'Airfield Data: aerodrome information.',
     ],
     'image': '14-global-met-menu.png',
     'image_pos': 'right',
     'image_caption': 'MET menu — weather products.'},

    {'kind': 'content',
     'title': 'Info — general & navaid information',
     'body': [
         'General Information, Misc Information, NOTAM display.',
         'Small QNH/TL window for quick reference.',
         'LFUNC Frequency window — list of frequencies in your area.',
         'Point label toggles: turn on/off Airport / Fix / NDB / VOR labels on the display.',
     ],
     'image': '15-global-info-menu.png',
     'image_pos': 'right',
     'image_caption': 'Info menu — general data plus point-label toggles.'},

    {'kind': 'content',
     'title': 'MSG — messages',
     'body': [
         'Quick access to ATC messaging windows and text notes.',
         'Used for inter-controller coordination via in-tool messaging.',
     ],
     'image': '16-global-msg-menu.png',
     'image_pos': 'right',
     'image_caption': 'MSG menu.'},

    {'kind': 'content',
     'title': 'STS — Status windows',
     'body': [
         'Plugin Status: TopSky internal state and supervisory.',
         'Safety Nets and TCT Status: see what safety nets are active.',
         'MTCD Status / FPCA Status / Divergence Detection — conflict probe configurations.',
         'CPDLC Default Status / FAST Configuration.',
         'Runway In Use / Runway Line Display: configure runway visualisation.',
         'Operational Load / Forecast windows: capacity vs demand visibility.',
     ],
     'image': '17-global-sts-menu.png',
     'image_pos': 'right',
     'image_caption': 'STS menu — status and configuration windows.'},

    # Sections 3 (Track Label menus), 5 (Windows), 6 (Aircraft Lists), 7 (Safety Nets),
    # 8 (Workflows), 9 (Resources) — placeholders for next session
    {'kind': 'divider', 'num': 4, 'name': '— remaining sections to be built next session —'},

    {'kind': 'content',
     'title': 'Coming in the next pass',
     'body': [
         'Section 4: Track Label menus (Callsign / AFL / CFL / RFL / AHDG / SSR / Transfer / etc.) — needs ~10 right-click captures on a live track.',
         'Section 5: Key Windows (Flight Plan, CARD, MTCD, AMS) — open each, capture, narrate.',
         'Section 6: Aircraft Lists (Sector, Holding, FAST, etc.) — capture each list configuration.',
         'Section 7: Safety Nets (AIW / APW / MSAW / STCA), Monitoring (CLAM / RAM), TCT.',
         'Section 8: Practical workflows — typical handoff sequence, CFL change, common shortcuts.',
         'Section 9: Resources — keyboard shortcut reference, doc pointers.',
     ]},
]


# --------------- assembly ---------------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for s in SLIDES:
        kind = s['kind']
        if kind == 'title':
            add_title_slide(prs, s['title'], s['subtitle'], s['footer'])
        elif kind == 'divider':
            add_section_divider(prs, s['num'], s['name'])
        elif kind == 'content':
            add_content_slide(
                prs,
                title=s['title'],
                body_lines=s.get('body'),
                image=s.get('image'),
                image_caption=s.get('image_caption'),
                image_pos=s.get('image_pos', 'none'),
            )

    prs.save(OUT_PATH)
    print(f'Wrote {len(SLIDES)} slides → {OUT_PATH}')


if __name__ == '__main__':
    main()
